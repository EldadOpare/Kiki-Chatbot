'''
This was our utilities script for ChromaDB operations and document processing for the Ghana chatbot.
We designed this script to handle various document types including PDFs, Word documents, PowerPoint presentations, 
CSV files, Excel files, and web scraping. All processed content was added to our ChromaDB vector database.

We implemented text cleaning, chunking strategies, and semantic text conversion to make tabular data 
more suitable for natural language processing and retrieval.

'''



#All Imports

import pymupdf
import requests
import pandas as pd
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation




def clean_text(text):
    """
    This function cleans and normalizes text by removing unwanted characters and formatting

    """
    
    #We replaced line breaks and carriage returns with spaces
    text = text.replace('\n', ' ').replace('\r', ' ')
    
    #We replaced tab characters with spaces
    text = text.replace('\t', ' ')
    
    #We removed form feed characters
    text = text.replace('\x0c', ' ')
    
    #We removed vertical tab characters
    text = text.replace('\x0b', ' ')

    #We normalized multiple spaces into single spaces
    text = ' '.join(text.split())
    
    #We removed leading and trailing whitespace
    text = text.strip()

    return text



def chunk_text(text, chunk_size=1000, overlap=200):
    """
    This function splits text into overlapping chunks for better context preservation
    """
    
    #We initialized our chunks list and starting position
    chunks = []
    start = 0
    text_length = len(text)
    
    #We processed the text in overlapping chunks
    while start < text_length:
        
        #We calculated the end position for this chunk
        end = start + chunk_size
        
        #We extracted the chunk from the text
        chunk = text[start:end]
        
        #We only added non-empty chunks to our list
        if chunk.strip():
            chunks.append(chunk.strip())
            
        #We moved the start position with overlap for context preservation
        start = start + chunk_size - overlap
    
    return chunks




def pdf_to_database(path, collection_name, original_filename=None):
    """
    This function extracts text from PDF files and adds to ChromaDB with page-level metadata
    """
    
    #We opened the PDF document using PyMuPDF
    doc = pymupdf.open(path)
    
    #We used the original filename if provided, otherwise extracted from path
    filename = original_filename if original_filename else path.split("/")[-1]  
    
    #We prepared lists to store all chunks, IDs, and metadata
    all_chunks = []
    all_ids = []
    all_metadata = []
    

    #We processed each page in the PDF document
    for page_number in range(len(doc)):
        page = doc[page_number]
        text = page.get_text()
    
        #We cleaned the extracted text
        text = clean_text(text)
        
        #We split the text into chunks for better retrieval
        chunks = chunk_text(text, chunk_size=1000, overlap=200)
        
        #We processed each chunk and created metadata
        for chunk_idx, chunk in enumerate(chunks):
            
            all_chunks.append(chunk)
            
            #We created unique IDs for each chunk
            all_ids.append(f"{filename}_p{page_number+1}_c{chunk_idx}")
            
            #We stored source, page, and chunk information in metadata
            all_metadata.append({
                "source": filename,
                "page": page_number + 1,
                "chunk": chunk_idx
            })
        
    #We added all chunks to the database if any were found
    if all_chunks:
        collection_name.upsert(
            documents=all_chunks,
            ids=all_ids,
            metadatas=all_metadata
        )

    #We closed the document to free memory
    doc.close()

    return


def scrapped_text_to_database(text,collection_name,source_name):
    """
    This function adds scraped or processed text to ChromaDB with appropriate chunking
    """
    
    #We prepared lists for chunks, IDs, and metadata
    all_chunks = []
    all_ids = []
    all_metadata = []

    #We cleaned the text before processing
    text = clean_text(text)
        
    #We split the text into manageable chunks
    chunks = chunk_text(text, chunk_size=1000, overlap=200)
        
    #We processed each chunk and created metadata
    for chunk_idx, chunk in enumerate(chunks):
            
        all_chunks.append(chunk)
            
        #We created unique IDs for each chunk
        all_ids.append(f"{source_name}_c{chunk_idx}")
            
        #We stored source and chunk information in metadata
        all_metadata.append({
            "source": source_name,
            "chunk": chunk_idx
        })
        
    #We added all chunks to the database if any were found
    if all_chunks:
        collection_name.upsert(
            documents=all_chunks,
            ids=all_ids,
            metadatas=all_metadata
        )

    return





def extract_pdf_text_from_bytes(pdf_content):
    """
    This function extracts text from PDF bytes using PyMuPDF for web-scraped PDF content
    """
    try:
        #We opened the PDF from bytes content
        doc = pymupdf.open(stream=pdf_content, filetype="pdf")
        
        #We extracted text from all pages
        text = ""
        for page in doc:
            text += page.get_text()
        
        #We closed the document to free memory
        doc.close()
        return text.strip()
    
    except Exception as e:
        return ""


def is_pdf_url(url):
    """
    This function checks if a URL pointed to a PDF file
    """
    return url.lower().endswith('.pdf')


def scrape_url(url, headers=None):
    """
    This function scrapes content from a single URL (HTML or PDF)
    """
    
    #We used more comprehensive headers to avoid blocking
    if headers is None:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'DNT': '1',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
        }
    
    try:
        #We made the HTTP request to get the content
        response = requests.get(url, headers=headers, timeout=15)
        
        # We handled different HTTP status codes appropriately
        if response.status_code == 403:
            print(f"Access forbidden (403) for {url} - site may have anti-bot protection")
            return None
        elif response.status_code == 404:
            print(f"Page not found (404) for {url}")
            return None
        elif response.status_code != 200:
            print(f"HTTP {response.status_code} error for {url}")
            return None
        
        #We checked if the content was a PDF
        if is_pdf_url(url) or 'application/pdf' in response.headers.get('Content-Type', ''):
            
            text = extract_pdf_text_from_bytes(response.content)
            return {
                'text': text,
                'type': 'pdf'
            }
        
        else:
            #For HTML content, we used BeautifulSoup to extract text
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # We removed script and style elements that don't contain useful content
            for script in soup(["script", "style"]):
                script.decompose()
            
            # We focused on paragraph tags for clean content
            paragraphs = soup.find_all('p')
            text = '\n\n'.join([p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)])
            
            # Only if no paragraphs found, we tried article or main content areas
            if not text.strip():
                articles = soup.find_all(['article', 'main'])
                for article in articles:
                    article_paragraphs = article.find_all('p')
                    if article_paragraphs:
                        text = '\n\n'.join([p.get_text(strip=True) for p in article_paragraphs if p.get_text(strip=True)])
                        break
            
            # We filtered out JavaScript requirement messages
            if 'enable javascript' in text.lower() and len(text) < 200:
                text = ""
            
            return {
                'text': text,
                'type': 'html'
            }
    
    except Exception as e:
        print(f"Error scraping {url}: {e}")
        return None


def scrape_url_to_database(url, collection_name):
    """
    This function scrapes a URL and adds the content to ChromaDB with appropriate source naming
    """
    
    #We scraped the URL to get its content
    result = scrape_url(url)
    
    if not result:
        print(f"Failed to scrape: {url}")
        return False

    if not result['text'].strip():
        print(f"No text content found at: {url}")
        return False
    
    #We extracted the text and content type
    text = result['text']
    content_type = result['type']
    
    if content_type == 'pdf':
        #For PDFs: we used the filename from URL as the source name
        filename = url.rstrip('/').split('/')[-1]
        source_name = filename if filename else url

        scrapped_text_to_database(
            text=text,
            collection_name=collection_name,
            source_name=source_name,
        )

    else:
        #For HTML: we used the URL as source name (removed trailing slash for consistency)
        source_name = url.rstrip('/') if url.rstrip('/') else url
        scrapped_text_to_database(
            text=text,
            collection_name=collection_name,
            source_name=source_name,
        )
    
    return True


def scrape_multiple_urls_to_database(urls, collection_name):
    """
    This function scrapes multiple URLs and adds all content to ChromaDB
    """

    #We processed each URL in the list
    for url in urls:
        result = scrape_url_to_database(url, collection_name)
        if result:
            print(f"Scraped and added to database: {url}")

    return


def docx_to_database(path, collection_name, original_filename=None):
    """
    This function extracts text from Word documents and adds to ChromaDB
    """

    #We opened the Word document
    doc = Document(path)
    filename = original_filename if original_filename else path.split("/")[-1]

    #We extracted all text from paragraphs
    full_text = ""
    for paragraph in doc.paragraphs:
        full_text += paragraph.text + "\n"

    #We also extracted text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text += cell.text + " "
            full_text += "\n"

    #We used the same processing pipeline as PDFs
    full_text = clean_text(full_text)

    #We added to database using our existing function
    scrapped_text_to_database(full_text, collection_name, filename)

    return


def pptx_to_database(path, collection_name, original_filename=None):
    """
    This function extracts text from PowerPoint presentations and adds to ChromaDB
    """

    #We opened the PowerPoint presentation
    prs = Presentation(path)
    filename = original_filename if original_filename else path.split("/")[-1]

    #We extracted text from all slides
    full_text = ""
    for slide_number, slide in enumerate(prs.slides, start=1):
        full_text += f"\n--- Slide {slide_number} ---\n"

        #We got text from all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"

    #We used the same processing pipeline as PDFs
    full_text = clean_text(full_text)

    #We added to database using our existing function
    scrapped_text_to_database(full_text, collection_name, filename)

    return


def dataframe_to_semantic_text(df, source_name, sheet_name=None):
    """
    This function converts DataFrames to natural language sentences for better NLP understanding.

    We designed this function to intelligently analyze data and create human-readable sentences
    that captured the relationships and information in the tabular data.

    """
    sentences = []

    #We added header context with sheet information if available
    sheet_context = f" from sheet '{sheet_name}'" if sheet_name else ""
    sentences.append(f"This is data from {source_name}{sheet_context}.")

    #We described the data structure
    num_rows, num_cols = df.shape
    column_names = ", ".join(df.columns.tolist())
    sentences.append(f"It contains {num_rows} records with {num_cols} attributes: {column_names}.")

    #We identified numeric columns for potential statistics
    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()

    #We added summary statistics for numeric columns if they existed
    if numeric_cols and len(df) > 0:
        stats_sentences = []
        for col in numeric_cols:  
            try:
                col_mean = df[col].mean()
                col_max = df[col].max()
                col_min = df[col].min()

                #We formatted numbers nicely for readability
                if col_mean > 1000:
                    stats_sentences.append(
                        f"The {col} values range from {col_min:,.0f} to {col_max:,.0f} with an average of {col_mean:,.0f}."
                    )
                else:
                    stats_sentences.append(
                        f"The {col} values range from {col_min:.2f} to {col_max:.2f} with an average of {col_mean:.2f}."
                    )
            except:
                pass

        if stats_sentences:
            sentences.extend(stats_sentences)

    #We converted rows to natural language sentences
    #We limited to reasonable number of rows to avoid token explosion
    max_rows_to_process = min(len(df), 400)

    for idx, row in df.head(max_rows_to_process).iterrows():
        
        #We built a natural sentence from the row
        row_parts = []

        for col, val in row.items():
            
            #We skipped NaN values
            if pd.isna(val):
                continue

            #We formatted the value appropriately
            if isinstance(val, (int, float)):
                if val > 1000:
                    formatted_val = f"{val:,.0f}"
                else:
                    formatted_val = f"{val}"
            else:
                formatted_val = str(val)

            row_parts.append(f"the {col} is {formatted_val}")

        #We created a sentence from the row
        if row_parts:
            if len(row_parts) == 1:
                sentence = f"There is a record where {row_parts[0]}."
            elif len(row_parts) == 2:
                sentence = f"There is a record where {row_parts[0]} and {row_parts[1]}."
            else:
                #We joined all but last with commas, last with "and"
                sentence = f"There is a record where {', '.join(row_parts[:-1])}, and {row_parts[-1]}."

            sentences.append(sentence)

    #We added a note if data was truncated
    if len(df) > max_rows_to_process:
        sentences.append(
            f"Note: This dataset contains {len(df)} total records, but only the first {max_rows_to_process} "
            f"records are detailed above for context and searchability."
        )

    #We joined all sentences with spaces
    return " ".join(sentences)


def csv_to_database(path, collection_name, original_filename=None):
    """
    This function extracts text from CSV files and adds to ChromaDB using semantic NLP-friendly format
    """

    #We read the CSV file into a DataFrame
    df = pd.read_csv(path)
    filename = original_filename if original_filename else path.split("/")[-1]

    #We converted the tabular data to semantic natural language
    full_text = dataframe_to_semantic_text(df, filename)

    #We used the same processing pipeline as PDFs
    full_text = clean_text(full_text)

    #We added to database using our existing function
    scrapped_text_to_database(full_text, collection_name, filename)

    return


def excel_to_database(path, collection_name, original_filename=None):
    """
    This function extracts text from Excel files (supported .xlsx and .xls) and adds to ChromaDB using semantic NLP-friendly format.
    """

    #We read the Excel file which automatically handled both .xlsx and .xls
    excel_file = pd.ExcelFile(path)
    filename = original_filename if original_filename else path.split("/")[-1]

    #We processed each sheet separately with semantic conversion
    all_text_parts = []

    #We added overall file context
    num_sheets = len(excel_file.sheet_names)
    sheet_names = ", ".join(excel_file.sheet_names)
    all_text_parts.append(
        f"This is an Excel file named {filename} containing {num_sheets} sheet(s): {sheet_names}."
    )

    #We processed each sheet in the Excel file
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(path, sheet_name=sheet_name)

        #We skipped empty sheets
        if df.empty:
            all_text_parts.append(f"The sheet '{sheet_name}' is empty.")
            continue

        #We converted each sheet to semantic text
        sheet_text = dataframe_to_semantic_text(df, filename, sheet_name=sheet_name)
        all_text_parts.append(sheet_text)

    #We combined all parts into one text
    full_text = " ".join(all_text_parts)

    #We used the same processing pipeline as PDFs
    full_text = clean_text(full_text)

    #We added to database using our existing function
    scrapped_text_to_database(full_text, collection_name, filename)

    return
